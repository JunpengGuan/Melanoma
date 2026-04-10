"""Generate an English, styled PowerPoint for the course project (Method 1 vs Method 2)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


def _blank_layout(prs: Presentation):
    for layout in prs.slide_layouts:
        if layout.name == "Blank":
            return layout
    return prs.slide_layouts[6]


def _set_fill(shape, rgb: tuple[int, int, int]) -> None:
    fill = getattr(shape, "fill", shape)
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb)


def _style_run(run, *, size_pt: int, bold: bool = False, color: RGBColor | None = None) -> None:
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.name = "Calibri"
    if color is not None:
        run.font.color.rgb = color


def add_title_slide(prs: Presentation, title: str, subtitle: str, footer: str) -> None:
    slide = prs.slides.add_slide(_blank_layout(prs))
    w, h = prs.slide_width, prs.slide_height
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, w, h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(23, 43, 77)
    bg.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, int(h * 0.42), int(w * 0.06), int(h * 0.14))
    _set_fill(accent, (56, 132, 209))
    accent.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.85), Inches(2.05), w - Inches(1.7), Inches(2.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.text = title
    if p.runs:
        _style_run(p.runs[0], size_pt=40, bold=True, color=RGBColor(255, 255, 255))

    p2 = tf.add_paragraph()
    p2.space_before = Pt(14)
    p2.alignment = PP_ALIGN.LEFT
    p2.text = subtitle
    if p2.runs:
        _style_run(p2.runs[0], size_pt=20, bold=False, color=RGBColor(198, 212, 235))

    foot = slide.shapes.add_textbox(Inches(0.85), h - Inches(1.0), w - Inches(1.7), Inches(0.55))
    ftf = foot.text_frame
    fp = ftf.paragraphs[0]
    fp.alignment = PP_ALIGN.LEFT
    fp.text = footer
    if fp.runs:
        _style_run(fp.runs[0], size_pt=14, bold=False, color=RGBColor(160, 174, 200))


def add_section_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(_blank_layout(prs))
    w, h = prs.slide_width, prs.slide_height

    body_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, w, h)
    _set_fill(body_bg, (252, 252, 253))
    body_bg.line.fill.background()

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, w, Inches(1.05))
    _set_fill(bar, (23, 43, 77))
    bar.line.fill.background()

    tb_title = slide.shapes.add_textbox(Inches(0.55), Inches(0.28), w - Inches(1.1), Inches(0.65))
    ttf = tb_title.text_frame
    ttf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tp = ttf.paragraphs[0]
    tp.alignment = PP_ALIGN.LEFT
    tp.text = title
    if tp.runs:
        _style_run(tp.runs[0], size_pt=26, bold=True, color=RGBColor(255, 255, 255))

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.05), Inches(0.12), h - Inches(1.05))
    _set_fill(accent, (56, 132, 209))
    accent.line.fill.background()

    body = slide.shapes.add_textbox(Inches(0.75), Inches(1.28), w - Inches(1.35), h - Inches(1.55))
    btf = body.text_frame
    btf.word_wrap = True
    body_rgb = RGBColor(45, 55, 72)
    for i, line in enumerate(bullets):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        p.text = line
        p.level = 0
        p.space_after = Pt(10)
        p.line_spacing = 1.15
        for r in p.runs:
            _style_run(r, size_pt=19, bold=False, color=body_rgb)


def add_table_slide(prs: Presentation, title: str, headers: list[str], rows: list[list[str]]) -> None:
    slide = prs.slides.add_slide(_blank_layout(prs))
    w, h = prs.slide_width, prs.slide_height

    body_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, w, h)
    _set_fill(body_bg, (252, 252, 253))
    body_bg.line.fill.background()

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, w, Inches(1.05))
    _set_fill(bar, (23, 43, 77))
    bar.line.fill.background()

    tb_title = slide.shapes.add_textbox(Inches(0.55), Inches(0.28), w - Inches(1.1), Inches(0.65))
    ttf = tb_title.text_frame
    ttf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tp = ttf.paragraphs[0]
    tp.text = title
    if tp.runs:
        _style_run(tp.runs[0], size_pt=26, bold=True, color=RGBColor(255, 255, 255))

    nrows = 1 + len(rows)
    ncols = len(headers)
    left = Inches(0.75)
    top = Inches(1.35)
    tbl_w = w - Inches(1.5)
    row_h = Inches(0.42)
    table = slide.shapes.add_table(nrows, ncols, left, top, tbl_w, row_h * nrows).table

    for j, htext in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = htext
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            _style_run(r, size_pt=15, bold=True, color=RGBColor(255, 255, 255))
        _set_fill(cell.fill, (56, 132, 209))

    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                _style_run(r, size_pt=15, bold=False, color=RGBColor(45, 55, 72))
            bg = (255, 255, 255) if i % 2 else (247, 250, 252)
            _set_fill(cell.fill, bg)


def main() -> None:
    root = Path(__file__).resolve().parents[1]
    # New filename avoids overwrite lock if the old deck is open in PowerPoint.
    out = root / "EE5271_Melanoma_Project_EN.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Melanoma vs Benign Skin Lesion Classification",
        "Comparing an end-to-end CNN with a segmentation + ABCD + tabular pipeline",
        "EE 5271 — Course Project · ISIC / ISBI 2016 Challenge Data",
    )

    add_section_slide(
        prs,
        "Objectives",
        [
            "Binary classification: benign vs malignant (melanoma) on dermoscopic crops.",
            "Method 1: strong black-box baseline (transfer learning from ImageNet).",
            "Method 2: interpretable pipeline — U-Net lesion mask → handcrafted features → LR / XGBoost.",
            "Report AUC-ROC, accuracy, sensitivity, specificity (threshold = 0.5) on a held-out test set.",
        ],
    )

    add_section_slide(
        prs,
        "Dataset & Protocol",
        [
            "Training: ~900 lesion-centered RGB images with image-level labels (Part 1 + training CSV).",
            "Validation: stratified hold-out (15%) for training monitoring; fixed random seed for reproducibility.",
            "Test: 379 images (Part 3B test) with binary labels — same split for both methods.",
            "Segmentation supervision: expert lesion masks (Part 1 ground truth) for U-Net training.",
        ],
    )

    add_section_slide(
        prs,
        "Method 1 — End-to-End CNN",
        [
            "Architecture: EfficientNet-B0, ImageNet-pretrained, fine-tuned head (1 logit).",
            "Input: full 224×224 RGB image — no explicit lesion mask.",
            "Optimization: AdamW + BCEWithLogitsLoss; class imbalance via weighted sampling and pos_weight.",
        ],
    )

    add_section_slide(
        prs,
        "Method 1 — Test Results (n = 379)",
        [
            "AUC-ROC ≈ 0.81 — strong ranking of malignant vs benign.",
            "Accuracy ≈ 0.76; sensitivity ≈ 0.65; specificity ≈ 0.78.",
            "Serves as a high-performance black-box reference for Method 2.",
        ],
    )

    add_section_slide(
        prs,
        "Method 2 — Interpretable Pipeline",
        [
            "Segmentation: U-Net (CNN encoder–decoder + skip connections); BCE + Dice loss.",
            "Features: 20-D ABCD-style descriptors (shape + color statistics inside predicted mask).",
            "Classifiers: Logistic Regression (linear, coefficient-level interpretability) and XGBoost (nonlinear baseline).",
        ],
    )

    add_section_slide(
        prs,
        "Method 2 — Test Results (n = 379)",
        [
            "Logistic Regression: AUC ≈ 0.74; sensitivity ≈ 0.72; specificity ≈ 0.67 — recommended primary tabular model.",
            "XGBoost: AUC ≈ 0.69; lower sensitivity (~0.44) — useful as a contrast / discussion of small-sample tree models.",
            "Trade-off vs Method 1: lower AUC but LR captures clinically motivated features; more false positives at 0.5 threshold.",
        ],
    )

    add_table_slide(
        prs,
        "Side-by-Side Summary (Test Set)",
        ["Approach", "AUC-ROC", "Accuracy", "Sensitivity", "Specificity"],
        [
            ["Method 1 — EfficientNet-B0", "≈ 0.81", "≈ 0.76", "≈ 0.65", "≈ 0.78"],
            ["Method 2 — ABCD + LR", "≈ 0.74", "≈ 0.68", "≈ 0.72", "≈ 0.67"],
            ["Method 2 — ABCD + XGB", "≈ 0.69", "≈ 0.71", "≈ 0.44", "≈ 0.77"],
        ],
    )

    add_section_slide(
        prs,
        "Conclusions & Next Steps",
        [
            "Both pipelines are end-to-end reproducible: training scripts, checkpoints, and shared test evaluation.",
            "Method 1 wins on AUC; Method 2 links predictions to explicit morphology and color descriptors.",
            "Future work: refine U-Net (Dice/IoU), enrich ABCD, threshold tuning for clinical cost asymmetry, k-fold CV.",
            "Thank you — questions welcome.",
        ],
    )

    prs.save(out)
    print(f"Wrote {out}")


if __name__ == "__main__":
    main()
